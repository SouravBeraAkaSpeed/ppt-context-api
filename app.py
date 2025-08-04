import requests
from flask import Flask, request, jsonify
from pptx import Presentation
from io import BytesIO
from flask_cors import CORS
import os
from urllib.parse import urlparse

# Import the Document class for DOCX files
from docx import Document

app = Flask(__name__)
CORS(app)

# New welcome route
@app.route('/', methods=['GET'])
def welcome():
    """
    A simple welcome route to confirm the API is running.
    """
    return jsonify({
        "message": "Welcome to the Document Text Extractor API!",
        "instructions": "Send a POST request to /extract-text with a JSON body containing 'file_uri' (e.g., {'file_uri': 'http://example.com/your_document.pptx'})"
    })

@app.route('/extract-text', methods=['POST'])
def extract_text():
    """
    Receives a file URI (PPTX or DOCX) via POST request, downloads the file,
    extracts text content, and returns it.
    """
    data = request.get_json()

    if not data or 'file_uri' not in data:
        return jsonify({"error": "Invalid input: 'file_uri' is required in the JSON body."}), 400

    file_uri = data['file_uri']

    try:
        # Parse the URL to get the file extension
        parsed_url = urlparse(file_uri)
        path = parsed_url.path
        _, file_extension = os.path.splitext(path)
        file_extension = file_extension.lower() # Ensure case-insensitivity

        if not file_extension:
            return jsonify({"error": "Could not determine file type from URI. Please ensure the URI contains a file extension (e.g., .pptx, .docx)."}), 400

        # Download the file content
        response = requests.get(file_uri, stream=True, verify=False)
        response.raise_for_status() # Raise an exception for bad status codes (4xx or 5xx)

        # Use BytesIO to treat the downloaded content as a file in memory
        file_content_in_memory = BytesIO(response.content)

        extracted_text = []
        formatted_text = ""

        if file_extension == '.pptx':
            # --- PPTX Extraction Logic ---
            prs = Presentation(file_content_in_memory)
            for slide_number, slide in enumerate(prs.slides):
                slide_text = f"--- Slide {slide_number + 1} ---\n"
                slide_content = []

                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text = shape.text.strip()
                        if text: # Only add non-empty text
                            slide_content.append(text)

                if slide_content:
                    slide_text += "\n".join(slide_content)
                else:
                    slide_text += "No visible text on this slide." # More descriptive
                extracted_text.append(slide_text)

            formatted_text = "\n\n".join(extracted_text)

        elif file_extension == '.docx':
            # --- DOCX Extraction Logic ---
            document = Document(file_content_in_memory)
            
            # Extract text from paragraphs
            # Join paragraphs with double newlines for better readability
            extracted_text_paragraphs = [p.text for p in document.paragraphs if p.text.strip()]
            formatted_text = "\n\n".join(extracted_text_paragraphs)
            
            # Optional: Add logic for tables if needed (more complex)
            # For simplicity, this example primarily focuses on paragraph text.

        else:
            return jsonify({"error": f"Unsupported file type: '{file_extension}'. Only .pptx and .docx files are supported."}), 400

        return jsonify({"success": True, "text_content": formatted_text})

    except requests.exceptions.RequestException as e:
        return jsonify({"error": f"Failed to download file from '{file_uri}': {e}"}), 500
    except Exception as e:
        # Catch specific errors from pptx/docx libraries for better debugging
        if "BadZipFile" in str(e) or "KeyError" in str(e) or "Package is not a Zip file" in str(e):
             return jsonify({"error": f"Failed to open/process the file. It might be corrupted or not a valid {file_extension} file. Error: {e}"}), 422
        return jsonify({"error": f"An unexpected error occurred while processing the file: {e}"}), 500


if __name__ == '__main__':
    app.run(host='0.0.0.0', port=5000, debug=True)